#!/usr/bin/env python3
"""
Script para analizar la estructura de un archivo PPTX
Extrae dimensiones, layouts, posiciones de elementos, etc.
"""

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    print("❌ python-pptx no está instalado")
    print("   Instalar con: pip install python-pptx")
    exit(1)

import os

def emu_to_inches(emu):
    """Convierte EMU (English Metric Units) a pulgadas"""
    return emu / 914400

def emu_to_cm(emu):
    """Convierte EMU a centímetros"""
    return emu_to_inches(emu) * 2.54

def analyze_ppt(filepath):
    print(f"📊 Analizando: {filepath}\n")
    
    if not os.path.exists(filepath):
        print(f"❌ Archivo no encontrado: {filepath}")
        return
    
    try:
        prs = Presentation(filepath)
        
        # Dimensiones del slide
        print("📐 DIMENSIONES DEL SLIDE:")
        width_inches = emu_to_inches(prs.slide_width)
        height_inches = emu_to_inches(prs.slide_height)
        width_cm = emu_to_cm(prs.slide_width)
        height_cm = emu_to_cm(prs.slide_height)
        
        print(f"   Ancho:  {width_inches:.2f}\" ({width_cm:.2f} cm)")
        print(f"   Alto:   {height_inches:.2f}\" ({height_cm:.2f} cm)")
        print(f"   Proporción: {width_inches/height_inches:.2f}:1\n")
        
        # Número de slides
        print(f"📄 TOTAL SLIDES: {len(prs.slides)}\n")
        
        # Layouts disponibles
        print(f"🎨 LAYOUTS DISPONIBLES: {len(prs.slide_layouts)}")
        for idx, layout in enumerate(prs.slide_layouts):
            print(f"   [{idx}] {layout.name}")
        print()
        
        # Analizar cada slide
        for slide_idx, slide in enumerate(prs.slides):
            print(f"\n{'='*60}")
            print(f"📄 SLIDE {slide_idx + 1}")
            print(f"   Layout: {slide.slide_layout.name}")
            print(f"{'='*60}")
            
            # Analizar shapes en el slide
            print(f"\n   🔲 ELEMENTOS EN EL SLIDE: {len(slide.shapes)}\n")
            
            for shape_idx, shape in enumerate(slide.shapes):
                print(f"   [{shape_idx}] {shape.shape_type} - {shape.name}")
                
                # Posición y dimensiones
                x_inches = emu_to_inches(shape.left)
                y_inches = emu_to_inches(shape.top)
                w_inches = emu_to_inches(shape.width)
                h_inches = emu_to_inches(shape.height)
                
                print(f"       Posición: X={x_inches:.2f}\", Y={y_inches:.2f}\"")
                print(f"       Tamaño:   W={w_inches:.2f}\", H={h_inches:.2f}\"")
                
                # Si es texto
                if hasattr(shape, 'text'):
                    text = shape.text.strip()
                    if text:
                        preview = text[:50] + "..." if len(text) > 50 else text
                        print(f"       Texto: \"{preview}\"")
                        
                        if hasattr(shape, 'text_frame'):
                            for para in shape.text_frame.paragraphs:
                                if para.runs:
                                    run = para.runs[0]
                                    if run.font.size:
                                        size_pt = run.font.size.pt
                                        print(f"       Fuente: {size_pt:.0f}pt", end="")
                                        if run.font.bold:
                                            print(" BOLD", end="")
                                        if run.font.color.rgb:
                                            print(f", Color: #{run.font.color.rgb}", end="")
                                        print()
                
                # Si es imagen
                if shape.shape_type == 13:  # PICTURE
                    print(f"       🖼️ IMAGEN detectada")
                    if hasattr(shape, 'image'):
                        print(f"       Formato: {shape.image.content_type}")
                
                print()
        
        print(f"\n{'='*60}")
        print("✅ Análisis completado")
        print(f"{'='*60}\n")
        
        # Resumen de configuración recomendada
        print("💡 CONFIGURACIÓN RECOMENDADA PARA PPTXGENJS:")
        print(f"\nconst pptx = new PptxGenJS();")
        print(f"// NO establecer layout - usar dimensiones por defecto 10x7.5")
        print(f"\nconst slide = pptx.addSlide();")
        print(f"\n// Para mantener proporciones del archivo original:")
        print(f"// Width: {width_inches:.2f}\" Height: {height_inches:.2f}\"")
        print(f"// (PptxGenJS default es 10x7.5, muy similar)\n")
        
    except Exception as e:
        print(f"❌ Error al analizar: {e}")
        import traceback
        traceback.print_exc()

if __name__ == "__main__":
    filepath = r"c:\Users\migduran\Documents\nuevo ooh\backend\REPORTE FACTURACIÓN BASE.pptx"
    analyze_ppt(filepath)
