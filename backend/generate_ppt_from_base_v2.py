#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Script para generar reporte PPT usando archivo base como plantilla
Estrategia: Llenar slide 0 → Duplicar → Actualizar → Repetir
"""

import sys
import json
import os
from shutil import copy2
from datetime import datetime

try:
    from pptx import Presentation
    from pptx.util import Inches, Cm, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    print("ERROR: python-pptx no instalado", file=sys.stderr)
    print("Instalar con: pip install python-pptx", file=sys.stderr)
    sys.exit(1)

def format_date_spanish(date_str):
    """Convierte yyyy-MM-dd a '4 de enero de 2026'"""
    months = ['enero','febrero','marzo','abril','mayo','junio',
              'julio','agosto','septiembre','octubre','noviembre','diciembre']
    try:
        parts = date_str.split('-')
        year, month, day = parts
        month_idx = int(month) - 1
        day_num = int(day)
        return f"{day_num} de {months[month_idx]} de {year}"
    except:
        return date_str

def convert_api_to_local_path(api_path):
    """Convierte /api/images/... a ruta local del filesystem"""
    if not api_path:
        return None
    
    script_dir = os.path.dirname(os.path.abspath(__file__))
    clean_path = api_path.replace('/api/images/', '', 1)
    full_path = os.path.join(script_dir, 'local-images', clean_path)
    normalized = os.path.normpath(full_path)
    
    return normalized

def duplicate_slide(prs, slide_index):
    """Duplica una diapositiva completa en la presentación"""
    from copy import deepcopy
    
    source_slide = prs.slides[slide_index]
    blank_layout = prs.slide_layouts[6]
    new_slide = prs.slides.add_slide(blank_layout)
    
    # Copiar SOLO textboxes y decoraciones (NO imágenes)
    for shape in source_slide.shapes:
        # Saltar imágenes (shape_type 13)
        if shape.shape_type == 13:
            print(f"[SKIP] No copiar imagen: {shape.name if hasattr(shape, 'name') else 'unknown'}")
            continue
        
        el = shape.element
        newel = deepcopy(el)
        new_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
    
    return new_slide

def update_slide_elements(slide, record, is_first=False):
    """
    Actualiza los elementos del slide con los datos del registro
    is_first: True para el primer registro (usa slide existente sin duplicación)
    """
    try:
        imagenes = record.get('imagenes', [])
        
        # Para la primera diapositiva, usar los índices conocidos
        if is_first:
            # [0] Dirección
            if len(slide.shapes) > 0 and hasattr(slide.shapes[0], 'text_frame'):
                tf = slide.shapes[0].text_frame
                tf.clear()
                p = tf.paragraphs[0]
                p.text = record['direccion'].upper()
                print(f"[0] Dirección: {record['direccion']}")
            
            # [1] Ciudad
            if len(slide.shapes) > 1 and hasattr(slide.shapes[1], 'text_frame'):
                tf = slide.shapes[1].text_frame
                tf.clear()
                p = tf.paragraphs[0]
                p.text = record['ciudad'].upper()
                print(f"[1] Ciudad: {record['ciudad']}")
            
            # [2] Vigencia
            if len(slide.shapes) > 2 and hasattr(slide.shapes[2], 'text_frame'):
                fecha_inicio = format_date_spanish(record['fechaInicio'])
                fecha_fin = format_date_spanish(record['fechaFin'])
                vigencia_text = f"Vigencia: {fecha_inicio} - {fecha_fin}"
                tf = slide.shapes[2].text_frame
                tf.clear()
                p = tf.paragraphs[0]
                p.text = vigencia_text
                print(f"[2] Vigencia actualizada")
            
            # [3] Proveedor (REF)
            if len(slide.shapes) > 3 and hasattr(slide.shapes[3], 'text_frame'):
                ref_text = f"REF: {record.get('proveedor', 'N/A')}"
                tf = slide.shapes[3].text_frame
                tf.clear()
                p = tf.paragraphs[0]
                p.text = ref_text
                if p.runs:
                    p.runs[0].font.color.rgb = RGBColor(0xBF, 0x90, 0x00)
                print(f"[3] Proveedor actualizado")
            
            # [6] Imagen 1 - Validar que existe
            if len(imagenes) > 0 and imagenes[0] and len(slide.shapes) > 6:
                img_path = convert_api_to_local_path(imagenes[0])
                print(f"[IMG] Validando imagen 1: {img_path}")
                if img_path and os.path.exists(img_path):
                    try:
                        placeholder = slide.shapes[6]
                        slide.shapes.add_picture(img_path, placeholder.left, placeholder.top, 
                                               width=placeholder.width, height=placeholder.height)
                        print(f"[IMG OK] Imagen 1 agregada correctamente")
                    except Exception as e:
                        print(f"[ERROR] Imagen 1: {e}", file=sys.stderr)
                else:
                    print(f"[WARN] Imagen 1 no encontrada: {imagenes[0]}", file=sys.stderr)
            
            # [7] Imagen 2 - Validar que existe
            if len(imagenes) > 1 and imagenes[1] and len(slide.shapes) > 7:
                img_path = convert_api_to_local_path(imagenes[1])
                print(f"[IMG] Validando imagen 2: {img_path}")
                if img_path and os.path.exists(img_path):
                    try:
                        placeholder = slide.shapes[7]
                        slide.shapes.add_picture(img_path, placeholder.left, placeholder.top,
                                               width=placeholder.width, height=placeholder.height)
                        print(f"[IMG OK] Imagen 2 agregada correctamente")
                    except Exception as e:
                        print(f"[ERROR] Imagen 2: {e}", file=sys.stderr)
                else:
                    print(f"[WARN] Imagen 2 no encontrada: {imagenes[1]}", file=sys.stderr)
            
            # [8] Imagen 3 - Validar que existe
            if len(imagenes) > 2 and imagenes[2] and len(slide.shapes) > 8:
                img_path = convert_api_to_local_path(imagenes[2])
                print(f"[IMG] Validando imagen 3: {img_path}")
                if img_path and os.path.exists(img_path):
                    try:
                        placeholder = slide.shapes[8]
                        slide.shapes.add_picture(img_path, placeholder.left, placeholder.top,
                                               width=placeholder.width, height=placeholder.height)
                        print(f"[IMG OK] Imagen 3 agregada correctamente")
                    except Exception as e:
                        print(f"[ERROR] Imagen 3: {e}", file=sys.stderr)
                else:
                    print(f"[WARN] Imagen 3 no encontrada: {imagenes[2]}", file=sys.stderr)
        
        else:
            # Para slides duplicadas, buscar por nombre y reemplazar
            # Las imágenes NO se copiaron, así que podemos agregar directamente
            
            for shape in slide.shapes:
                shape_name = shape.name if hasattr(shape, 'name') else ''
                
                # Actualizar textos por nombre
                if shape_name == 'Título 4' and hasattr(shape, 'text_frame'):
                    tf = shape.text_frame
                    tf.clear()
                    p = tf.paragraphs[0]
                    p.text = record['direccion'].upper()
                    print(f"[TXT] Título: {record['direccion']}")
                
                elif shape_name == 'Marcador de texto 5' and hasattr(shape, 'text_frame'):
                    tf = shape.text_frame
                    tf.clear()
                    p = tf.paragraphs[0]
                    p.text = record['ciudad'].upper()
                    print(f"[TXT] Ciudad: {record['ciudad']}")
                
                elif shape_name == 'Marcador de texto 8' and hasattr(shape, 'text_frame'):
                    fecha_inicio = format_date_spanish(record['fechaInicio'])
                    fecha_fin = format_date_spanish(record['fechaFin'])
                    vigencia_text = f"Vigencia: {fecha_inicio} - {fecha_fin}"
                    tf = shape.text_frame
                    tf.clear()
                    p = tf.paragraphs[0]
                    p.text = vigencia_text
                    print(f"[TXT] Vigencia")
                
                elif shape_name == 'Marcador de texto 7' and hasattr(shape, 'text_frame'):
                    ref_text = f"REF: {record.get('proveedor', 'N/A')}"
                    tf = shape.text_frame
                    tf.clear()
                    p = tf.paragraphs[0]
                    p.text = ref_text
                    if p.runs:
                        p.runs[0].font.color.rgb = RGBColor(0xBF, 0x90, 0x00)
                    print(f"[TXT] Proveedor")
            
            # Agregar imágenes nuevas (no hay viejas que reemplazar)
            if len(imagenes) > 0 and imagenes[0]:
                img_path = convert_api_to_local_path(imagenes[0])
                print(f"[IMG] Validando imagen 1 en duplicado: {img_path}")
                if img_path and os.path.exists(img_path):
                    try:
                        # Obtener posición del placeholder original
                        placeholder_1 = None
                        for shape in prs.slides[len(prs.slides) - 2].shapes:
                            if hasattr(shape, 'name') and shape.name == 'Marcador de posición de imagen 2':
                                placeholder_1 = shape
                                break
                        
                        if placeholder_1:
                            slide.shapes.add_picture(img_path, placeholder_1.left, placeholder_1.top,
                                                   width=placeholder_1.width, height=placeholder_1.height)
                            print(f"[IMG OK] Imagen 1 agregada en duplicado")
                    except Exception as e:
                        print(f"[ERROR] Imagen 1 duplicado: {e}", file=sys.stderr)
                else:
                    print(f"[WARN] Imagen 1 no encontrada: {imagenes[0]}", file=sys.stderr)
            
            if len(imagenes) > 1 and imagenes[1]:
                img_path = convert_api_to_local_path(imagenes[1])
                print(f"[IMG] Validando imagen 2 en duplicado: {img_path}")
                if img_path and os.path.exists(img_path):
                    try:
                        placeholder_2 = None
                        for shape in prs.slides[len(prs.slides) - 2].shapes:
                            if hasattr(shape, 'name') and shape.name == 'Marcador de posición de imagen 9':
                                placeholder_2 = shape
                                break
                        
                        if placeholder_2:
                            slide.shapes.add_picture(img_path, placeholder_2.left, placeholder_2.top,
                                                   width=placeholder_2.width, height=placeholder_2.height)
                            print(f"[IMG OK] Imagen 2 agregada en duplicado")
                    except Exception as e:
                        print(f"[ERROR] Imagen 2 duplicado: {e}", file=sys.stderr)
                else:
                    print(f"[WARN] Imagen 2 no encontrada: {imagenes[1]}", file=sys.stderr)
            
            if len(imagenes) > 2 and imagenes[2]:
                img_path = convert_api_to_local_path(imagenes[2])
                print(f"[IMG] Validando imagen 3 en duplicado: {img_path}")
                if img_path and os.path.exists(img_path):
                    try:
                        placeholder_3 = None
                        for shape in prs.slides[len(prs.slides) - 2].shapes:
                            if hasattr(shape, 'name') and shape.name == 'Marcador de posición de imagen 11':
                                placeholder_3 = shape
                                break
                        
                        if placeholder_3:
                            slide.shapes.add_picture(img_path, placeholder_3.left, placeholder_3.top,
                                                   width=placeholder_3.width, height=placeholder_3.height)
                            print(f"[IMG OK] Imagen 3 agregada en duplicado")
                    except Exception as e:
                        print(f"[ERROR] Imagen 3 duplicado: {e}", file=sys.stderr)
                else:
                    print(f"[WARN] Imagen 3 no encontrada: {imagenes[2]}", file=sys.stderr)
    
    except Exception as e:
        print(f"[ERROR] Error al actualizar elementos: {e}", file=sys.stderr)

def main():
    if len(sys.argv) < 2:
        print("Uso: python generate_ppt_from_base_v2.py <data.json>", file=sys.stderr)
        sys.exit(1)
    
    data_file = sys.argv[1]
    
    # Leer datos
    try:
        with open(data_file, 'r', encoding='utf-8') as f:
            data = json.load(f)
    except Exception as e:
        print(f"[ERROR] No se pudo leer {data_file}: {e}", file=sys.stderr)
        sys.exit(1)
    
    base_file = data.get('base_file')
    output_file = data.get('output_file')
    records = data.get('records', [])
    month = data.get('month', '')
    
    if not base_file or not output_file:
        print("[ERROR] Faltan parametros: base_file o output_file", file=sys.stderr)
        sys.exit(1)
    
    print(f"[PPT] Generando PPT con {len(records)} registros para {month}")
    print(f"   Base: {base_file}")
    print(f"   Output: {output_file}")
    
    # Verificar que existe el archivo base
    if not os.path.exists(base_file):
        print(f"[ERROR] Archivo base no encontrado: {base_file}", file=sys.stderr)
        sys.exit(1)
    
    # Cargar presentacion base
    try:
        prs = Presentation(base_file)
        print(f"[OK] Archivo base cargado ({len(prs.slides)} slides existentes)")
    except Exception as e:
        print(f"[ERROR] No se pudo cargar el archivo base: {e}", file=sys.stderr)
        sys.exit(1)
    
    # Procesar cada registro
    for idx, record in enumerate(records):
        print(f"\n{'='*60}")
        if idx == 0:
            # Primer registro: usar slide 0 existente
            slide = prs.slides[0]
            print(f"[SLIDE 1] Rellenando slide base con primer registro")
            update_slide_elements(slide, record, is_first=True)
        else:
            # Siguientes registros: duplicar el slide anterior y actualizar
            print(f"[SLIDE {idx + 1}] Duplicando slide anterior")
            slide = duplicate_slide(prs, len(prs.slides) - 1)
            update_slide_elements(slide, record, is_first=False)
        
        print(f"[OK] Slide {idx + 1}: {record['direccion']} ({record['ciudad']})")
    
    # Guardar
    try:
        prs.save(output_file)
        print(f"\n{'='*60}")
        print(f"[OK] Archivo generado: {output_file}")
        print(f"[OK] Total slides: {len(prs.slides)}")
        print(f"{'='*60}\n")
    except Exception as e:
        print(f"[ERROR] No se pudo guardar el archivo: {e}", file=sys.stderr)
        sys.exit(1)

if __name__ == '__main__':
    main()
