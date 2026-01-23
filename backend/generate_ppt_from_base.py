#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Script para generar reporte PPT usando archivo base como plantilla
Usa python-pptx que SI puede leer y modificar archivos PPTX existentes
SIN EMOJIS - Compatible con encoding cp1252 de Windows
"""

import sys
import json
import os
from shutil import copy2
from datetime import datetime

try:
    from pptx import Presentation
    from pptx.util import Inches, Cm, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    print("ERROR: python-pptx no instalado", file=sys.stderr)
    print("Instalar con: pip install python-pptx", file=sys.stderr)
    sys.exit(1)

# Colores ABI
COLOR_ROJO = RGBColor(204, 0, 0)      # #CC0000
COLOR_ORO = RGBColor(212, 165, 116)   # #D4A574
COLOR_AZUL = RGBColor(0, 51, 102)     # #003366
COLOR_GRIS = RGBColor(102, 102, 102)  # #666666

def format_date_spanish(date_str):
    """Convierte yyyy-MM-dd a '4 de enero de 2026'"""
    months = ['enero','febrero','marzo','abril','mayo','junio',
              'julio','agosto','septiembre','octubre','noviembre','diciembre']
    try:
        parts = date_str.split('-')
        year, month, day = parts
        month_idx = int(month) - 1
        day_num = int(day)
        return f"{day_num} de {months[month_idx]} de {year}"
    except:
        return date_str

def duplicate_slide(prs, slide_index):
    """Duplica una diapositiva completa en la presentación"""
    from copy import deepcopy
    
    source_slide = prs.slides[slide_index]
    blank_layout = prs.slide_layouts[6]  # Diseño en blanco
    new_slide = prs.slides.add_slide(blank_layout)
    
    # Copiar todas las shapes de la diapositiva origen EXCEPTO las imágenes
    for shape in source_slide.shapes:
        # Saltar imágenes (shape_type 13)
        if shape.shape_type == 13:
            continue
        
        el = shape.element
        newel = deepcopy(el)
        new_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
    
    return new_slide

def fill_slide_elements(slide, record):
    """
    Llena los elementos del slide con los datos del registro
    
    Índices esperados:
    [0] Título 4 → Dirección
    [1] Marcador de texto 5 → Ciudad
    [2] Marcador de texto 8 → Vigencia
    [3] Marcador de texto 7 → Proveedor (REF)
    [6] Marcador de posición de imagen 2 → Imagen 1
    [7] Marcador de posición de imagen 9 → Imagen 2
    [8] Marcador de posición de imagen 11 → Imagen 3
    """
    try:
        # [0] Dirección
        if len(slide.shapes) > 0 and hasattr(slide.shapes[0], 'text_frame'):
            tf = slide.shapes[0].text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = record['direccion'].upper()
            print(f"[0] Dirección: {record['direccion']}")
        
        # [1] Ciudad
        if len(slide.shapes) > 1 and hasattr(slide.shapes[1], 'text_frame'):
            tf = slide.shapes[1].text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = record['ciudad'].upper()
            print(f"[1] Ciudad: {record['ciudad']}")
        
        # [2] Vigencia
        if len(slide.shapes) > 2 and hasattr(slide.shapes[2], 'text_frame'):
            fecha_inicio = format_date_spanish(record['fechaInicio'])
            fecha_fin = format_date_spanish(record['fechaFin'])
            vigencia_text = f"Vigencia: {fecha_inicio} - {fecha_fin}"
            tf = slide.shapes[2].text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = vigencia_text
            print(f"[2] Vigencia: {vigencia_text}")
        
        # [3] Proveedor (REF) - Color #BF9000
        if len(slide.shapes) > 3 and hasattr(slide.shapes[3], 'text_frame'):
            ref_text = f"REF: {record.get('proveedor', 'N/A')}"
            tf = slide.shapes[3].text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = ref_text
            # Aplicar color #BF9000 al texto
            if p.runs:
                p.runs[0].font.color.rgb = RGBColor(0xBF, 0x90, 0x00)
            print(f"[3] Proveedor: {ref_text}")
        
        # Llenar imágenes
        imagenes = record.get('imagenes', [])
        
        # Primero, eliminar cualquier imagen existente (de la duplicación)
        # Iterar hacia atrás para no afectar los índices
        for idx in range(len(slide.shapes) - 1, -1, -1):
            if slide.shapes[idx].shape_type == 13:  # 13 = PICTURE
                sp = slide.shapes[idx]._element
                sp.getparent().remove(sp)
                print(f"[IMG] Removida imagen antigua en índice {idx}")
        
        # Ahora agregar las nuevas imágenes en las posiciones correctas
        # Reconstruir referencias a los placeholders (ya que borramos imágenes)
        placeholder_1 = None
        placeholder_2 = None
        placeholder_3 = None
        
        # Buscar los placeholders de imagen por nombre o tamaño
        for idx, shape in enumerate(slide.shapes):
            if hasattr(shape, 'name'):
                if 'imagen 2' in shape.name.lower() or 'posición de imagen 2' in shape.name.lower():
                    placeholder_1 = shape
                elif 'imagen 9' in shape.name.lower() or 'posición de imagen 9' in shape.name.lower():
                    placeholder_2 = shape
                elif 'imagen 11' in shape.name.lower() or 'posición de imagen 11' in shape.name.lower():
                    placeholder_3 = shape
        
        # [6] Imagen 1 (grande)
        if len(imagenes) > 0 and imagenes[0]:
            img_path = convert_api_to_local_path(imagenes[0])
            print(f"[DEBUG] Imagen 1: {img_path}")
            print(f"[DEBUG] Existe: {os.path.exists(img_path) if img_path else 'N/A'}")
            if img_path and os.path.exists(img_path):
                try:
                    if placeholder_1:
                        slide.shapes.add_picture(
                            img_path,
                            placeholder_1.left,
                            placeholder_1.top,
                            width=placeholder_1.width,
                            height=placeholder_1.height
                        )
                        print(f"[IMG OK] Imagen 1 agregada: {os.path.basename(img_path)}")
                    else:
                        print(f"[WARN] Placeholder 1 no encontrado")
                except Exception as e:
                    print(f"[ERROR] Imagen 1: {e}", file=sys.stderr)
            else:
                print(f"[WARN] Imagen 1 no encontrada: {imagenes[0]}", file=sys.stderr)
        
        # [7] Imagen 2 (pequeña arriba derecha)
        if len(imagenes) > 1 and imagenes[1]:
            img_path = convert_api_to_local_path(imagenes[1])
            print(f"[DEBUG] Imagen 2: {img_path}")
            if img_path and os.path.exists(img_path):
                try:
                    if placeholder_2:
                        slide.shapes.add_picture(
                            img_path,
                            placeholder_2.left,
                            placeholder_2.top,
                            width=placeholder_2.width,
                            height=placeholder_2.height
                        )
                        print(f"[IMG OK] Imagen 2 agregada: {os.path.basename(img_path)}")
                    else:
                        print(f"[WARN] Placeholder 2 no encontrado")
                except Exception as e:
                    print(f"[ERROR] Imagen 2: {e}", file=sys.stderr)
            else:
                print(f"[WARN] Imagen 2 no encontrada: {imagenes[1]}", file=sys.stderr)
        
        # [8] Imagen 3 (pequeña abajo derecha)
        if len(imagenes) > 2 and imagenes[2]:
            img_path = convert_api_to_local_path(imagenes[2])
            print(f"[DEBUG] Imagen 3: {img_path}")
            if img_path and os.path.exists(img_path):
                try:
                    if placeholder_3:
                        slide.shapes.add_picture(
                            img_path,
                            placeholder_3.left,
                            placeholder_3.top,
                            width=placeholder_3.width,
                            height=placeholder_3.height
                        )
                        print(f"[IMG OK] Imagen 3 agregada: {os.path.basename(img_path)}")
                    else:
                        print(f"[WARN] Placeholder 3 no encontrado")
                except Exception as e:
                    print(f"[ERROR] Imagen 3: {e}", file=sys.stderr)
            else:
                print(f"[WARN] Imagen 3 no encontrada: {imagenes[2]}", file=sys.stderr)
    
    except Exception as e:
        print(f"[ERROR] Error al llenar elementos: {e}", file=sys.stderr)

def add_vaya_slide(prs, record, slide_number):
    """
    Agrega un slide duplicando el slide anterior y rellenando con datos del registro
    """
    if slide_number == 1:
        # Primer registro: modificar el slide 0 (template)
        slide = prs.slides[0]
    else:
        # Siguientes registros: duplicar el slide anterior
        slide = duplicate_slide(prs, len(prs.slides) - 1)
    
    # Llenar elementos del slide con datos del registro
    fill_slide_elements(slide, record)
    
    print(f"[OK] Slide {slide_number}: {record['direccion']} ({record['ciudad']})")

def convert_api_to_local_path(api_path):
    """Convierte /api/images/... a ruta local del filesystem"""
    if not api_path:
        return None
    
    # Script está en backend/, imágenes están en backend/images/
    # api_path es: /api/images/MARCA/CAMPANA/YYYY-MM/filename.jpg
    
    script_dir = os.path.dirname(os.path.abspath(__file__))  # backend/
    
    # Remover /api/images/ del inicio
    clean_path = api_path.replace('/api/images/', '', 1)  # Resultado: MARCA/CAMPANA/...
    
    # Construir ruta en local-images
    full_path = os.path.join(script_dir, 'local-images', clean_path)
    normalized = os.path.normpath(full_path)
    
    return normalized

def main():
    if len(sys.argv) < 2:
        print("Uso: python generate_ppt_from_base.py <data.json>", file=sys.stderr)
        sys.exit(1)
    
    data_file = sys.argv[1]
    
    # Leer datos
    try:
        with open(data_file, 'r', encoding='utf-8') as f:
            data = json.load(f)
    except Exception as e:
        print(f"[ERROR] No se pudo leer {data_file}: {e}", file=sys.stderr)
        sys.exit(1)
    
    base_file = data.get('base_file')
    output_file = data.get('output_file')
    records = data.get('records', [])
    month = data.get('month', '')
    
    if not base_file or not output_file:
        print("[ERROR] Faltan parametros: base_file o output_file", file=sys.stderr)
        sys.exit(1)
    
    print(f"[PPT] Generando PPT con {len(records)} registros para {month}")
    print(f"   Base: {base_file}")
    print(f"   Output: {output_file}")
    
    # Verificar que existe el archivo base
    if not os.path.exists(base_file):
        print(f"[ERROR] Archivo base no encontrado: {base_file}", file=sys.stderr)
        sys.exit(1)
    
    # Cargar presentacion base
    try:
        prs = Presentation(base_file)
        print(f"[OK] Archivo base cargado ({len(prs.slides)} slides existentes)")
    except Exception as e:
        print(f"[ERROR] No se pudo cargar el archivo base: {e}", file=sys.stderr)
        sys.exit(1)
    
    # Agregar slides de VAYAS
    for idx, record in enumerate(records):
        try:
            add_vaya_slide(prs, record, idx + 1)
        except Exception as e:
            print(f"[ERROR] Slide {idx + 1}: {e}", file=sys.stderr)
    
    # Guardar
    try:
        prs.save(output_file)
        print(f"[OK] Archivo generado: {output_file}")
        print(f"[OK] Total slides: {len(prs.slides)}")
    except Exception as e:
        print(f"[ERROR] No se pudo guardar el archivo: {e}", file=sys.stderr)
        sys.exit(1)

if __name__ == '__main__':
    main()
