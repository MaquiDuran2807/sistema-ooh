#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Script para generar reporte PPT usando archivo base
Estrategia: Cargar base → Llenar cada slide disponible → Guardar
"""

import sys
import json
import os
from datetime import datetime

try:
    from pptx import Presentation
    from pptx.util import Inches, Cm, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    print("ERROR: python-pptx no instalado", file=sys.stderr)
    print("Instalar con: pip install python-pptx", file=sys.stderr)
    sys.exit(1)

def format_date_spanish(date_str):
    """Convierte yyyy-MM-dd a '4 de enero de 2026'"""
    months = ['enero','febrero','marzo','abril','mayo','junio',
              'julio','agosto','septiembre','octubre','noviembre','diciembre']
    try:
        parts = date_str.split('-')
        year, month, day = parts
        month_idx = int(month) - 1
        day_num = int(day)
        return f"{day_num} de {months[month_idx]} de {year}"
    except:
        return date_str

def convert_api_to_local_path(api_path):
    """Convierte /api/images/... a ruta local del filesystem"""
    if not api_path:
        return None
    
    # Si ya es ruta absoluta, devolver normalizada
    if os.path.isabs(api_path):
        return os.path.normpath(api_path)

    script_dir = os.path.dirname(os.path.abspath(__file__))
    clean_path = api_path.replace('/api/images/', '', 1)
    full_path = os.path.join(script_dir, 'local-images', clean_path)
    normalized = os.path.normpath(full_path)

    # Si existe, devolver directo
    if os.path.exists(normalized):
        return normalized

    # Fallback: buscar por nombre de archivo en todo local-images (por si la ruta viene mal en el CSV)
    filename = os.path.basename(clean_path)
    local_root = os.path.join(script_dir, 'local-images')
    for root, _, files in os.walk(local_root):
        if filename in files:
            candidate = os.path.join(root, filename)
            print(f"[FALLBACK] Imagen encontrada por búsqueda: {candidate}")
            return candidate

    return normalized

def find_placeholder(slide, name, fallback_index=None):
    """Busca un shape por nombre; si no lo encuentra usa índice de respaldo"""
    for shape in slide.shapes:
        if getattr(shape, "name", "") == name:
            return shape
    if fallback_index is not None and len(slide.shapes) > fallback_index:
        return slide.shapes[fallback_index]
    return None


def validate_image_placeholders(slide):
    """Loguea posiciones/tamaños de placeholders de imagen del slide base"""
    names = [
        "Marcador de posición de imagen 2",
        "Marcador de posición de imagen 9",
        "Marcador de posición de imagen 11",
    ]
    for idx, name in enumerate(names, start=1):
        ph = find_placeholder(slide, name, fallback_index=5 + idx)
        if ph:
            print(
                f"[CHECK] Placeholder img{idx}: {name} -> left={ph.left}, top={ph.top}, width={ph.width}, height={ph.height}"
            )
        else:
            print(f"[WARN] No se encontró placeholder para img{idx}: {name}")


def fill_slide(slide, record):
    """
    Llena un slide con los datos del registro
    Asume que el slide tiene los 9 elementos en las posiciones correctas
    """
    imagenes = record.get('imagenes', [])
    
    # [0] Dirección
    if len(slide.shapes) > 0 and hasattr(slide.shapes[0], 'text_frame'):
        tf = slide.shapes[0].text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = record['direccion'].upper()
        print(f"  [TXT] Dirección: {record['direccion']}")
    
    # [1] Ciudad
    if len(slide.shapes) > 1 and hasattr(slide.shapes[1], 'text_frame'):
        tf = slide.shapes[1].text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = record['ciudad'].upper()
        print(f"  [TXT] Ciudad: {record['ciudad']}")
    
    # [2] Vigencia
    if len(slide.shapes) > 2 and hasattr(slide.shapes[2], 'text_frame'):
        fecha_inicio = format_date_spanish(record['fechaInicio'])
        fecha_fin = format_date_spanish(record['fechaFin'])
        vigencia_text = f"Vigencia: {fecha_inicio} - {fecha_fin}"
        tf = slide.shapes[2].text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = vigencia_text
        print(f"  [TXT] Vigencia actualizada")
    
    # [3] Proveedor (REF)
    if len(slide.shapes) > 3 and hasattr(slide.shapes[3], 'text_frame'):
        ref_text = f"REF: {record.get('proveedor', 'N/A')}"
        tf = slide.shapes[3].text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = ref_text
        if p.runs:
            p.runs[0].font.color.rgb = RGBColor(0xBF, 0x90, 0x00)
        print(f"  [TXT] Proveedor: {record.get('proveedor', 'N/A')}")
    
    # Preparar placeholders (se obtienen antes de eliminar para evitar desalineación)
    ph1 = find_placeholder(slide, "Marcador de posición de imagen 2", fallback_index=6)
    ph2 = find_placeholder(slide, "Marcador de posición de imagen 9", fallback_index=7)
    ph3 = find_placeholder(slide, "Marcador de posición de imagen 11", fallback_index=8)

    # [6] Imagen 1
    if len(imagenes) > 0 and imagenes[0] and ph1:
        img_path = convert_api_to_local_path(imagenes[0])
        if img_path and os.path.exists(img_path):
            try:
                left, top, width, height = ph1.left, ph1.top, ph1.width, ph1.height
                sp = ph1.element
                sp.getparent().remove(sp)
                slide.shapes.add_picture(img_path, left, top, width=width, height=height)
                print(f"  [IMG] Imagen 1 agregada")
            except Exception as e:
                print(f"  [ERROR] Imagen 1: {e}", file=sys.stderr)
        else:
            print(f"  [WARN] Imagen 1 no encontrada: {imagenes[0]}", file=sys.stderr)
    elif len(imagenes) > 0 and imagenes[0]:
        print(f"  [WARN] No hay placeholder para Imagen 1")

    # [7] Imagen 2
    if len(imagenes) > 1 and imagenes[1] and ph2:
        img_path = convert_api_to_local_path(imagenes[1])
        if img_path and os.path.exists(img_path):
            try:
                left, top, width, height = ph2.left, ph2.top, ph2.width, ph2.height
                sp = ph2.element
                sp.getparent().remove(sp)
                slide.shapes.add_picture(img_path, left, top, width=width, height=height)
                print(f"  [IMG] Imagen 2 agregada")
            except Exception as e:
                print(f"  [ERROR] Imagen 2: {e}", file=sys.stderr)
        else:
            print(f"  [WARN] Imagen 2 no encontrada: {imagenes[1]}", file=sys.stderr)
    elif len(imagenes) > 1 and imagenes[1]:
        print(f"  [WARN] No hay placeholder para Imagen 2")

    # [8] Imagen 3
    if len(imagenes) > 2 and imagenes[2] and ph3:
        img_path = convert_api_to_local_path(imagenes[2])
        if img_path and os.path.exists(img_path):
            try:
                left, top, width, height = ph3.left, ph3.top, ph3.width, ph3.height
                sp = ph3.element
                sp.getparent().remove(sp)
                slide.shapes.add_picture(img_path, left, top, width=width, height=height)
                print(f"  [IMG] Imagen 3 agregada")
            except Exception as e:
                print(f"  [ERROR] Imagen 3: {e}", file=sys.stderr)
        else:
            print(f"  [WARN] Imagen 3 no encontrada: {imagenes[2]}", file=sys.stderr)
    elif len(imagenes) > 2 and imagenes[2]:
        print(f"  [WARN] No hay placeholder para Imagen 3")

def remove_blank_slides(prs, filled_count):
    """
    Elimina todos los slides posteriores al último slide lleno
    """
    total_slides = len(prs.slides)
    slides_to_remove = total_slides - filled_count
    
    if slides_to_remove <= 0:
        return
    
    print(f"[CLEAN] Eliminando {slides_to_remove} slides en blanco...")
    
    # Eliminar de atrás hacia adelante para no afectar los índices
    for i in range(total_slides - 1, filled_count - 1, -1):
        rId = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[i]
    
    print(f"[CLEAN] {slides_to_remove} slides eliminados")

def main():
    if len(sys.argv) < 2:
        print("Uso: python generate_ppt_from_base_v3.py <data.json>", file=sys.stderr)
        sys.exit(1)
    
    data_file = sys.argv[1]
    
    # Leer datos
    try:
        with open(data_file, 'r', encoding='utf-8') as f:
            data = json.load(f)
    except Exception as e:
        print(f"[ERROR] No se pudo leer {data_file}: {e}", file=sys.stderr)
        sys.exit(1)
    
    base_file = data.get('base_file')
    output_file = data.get('output_file')
    records = data.get('records', [])
    month = data.get('month', '')
    
    if not base_file or not output_file:
        print("[ERROR] Faltan parametros: base_file o output_file", file=sys.stderr)
        sys.exit(1)
    
    print(f"[PPT] Generando PPT con {len(records)} registros para {month}")
    print(f"   Base: {base_file}")
    print(f"   Output: {output_file}")
    
    # Verificar que existe el archivo base
    if not os.path.exists(base_file):
        print(f"[ERROR] Archivo base no encontrado: {base_file}", file=sys.stderr)
        sys.exit(1)
    
    # Cargar presentacion base
    prs = Presentation(base_file)
    print(f"\n[BASE] Cargada con {len(prs.slides)} slides disponibles")
    
    # Validar posiciones de placeholders con el slide base (primero)
    if len(prs.slides) > 0:
        validate_image_placeholders(prs.slides[0])

    # Verificar que hay suficientes slides
    if len(prs.slides) < len(records):
        print(f"[WARN] Base tiene {len(prs.slides)} slides pero hay {len(records)} registros")
        print(f"[WARN] Solo se llenaran los primeros {len(prs.slides)} registros")
        records = records[:len(prs.slides)]
    
    # Llenar cada slide con su registro
    print(f"\n{'='*60}")
    print("[LLENANDO SLIDES]")
    print(f"{'='*60}\n")
    
    for idx, record in enumerate(records):
        print(f"[SLIDE {idx + 1}] {record['direccion']} ({record['ciudad']})")
        try:
            slide = prs.slides[idx]
            fill_slide(slide, record)
        except Exception as e:
            print(f"[ERROR] No se pudo llenar slide {idx + 1}: {e}", file=sys.stderr)
    
    # Guardar resultado final
    print(f"\n{'='*60}")
    print(f"[LIMPIANDO]")
    print(f"{'='*60}\n")
    
    # Eliminar slides en blanco
    remove_blank_slides(prs, len(records))
    
    print(f"\n{'='*60}")
    print(f"[GUARDANDO]")
    print(f"{'='*60}\n")
    
    try:
        prs.save(output_file)
        print(f"[OK] Presentación guardada: {output_file}")
        print(f"[OK] Total de slides finales: {len(prs.slides)}")
    except Exception as e:
        print(f"[ERROR] No se pudo guardar: {e}", file=sys.stderr)
        sys.exit(1)
    
    print(f"\n{'='*60}")
    print(f"[OK] Proceso completado exitosamente")
    print(f"{'='*60}\n")

if __name__ == '__main__':
    main()
