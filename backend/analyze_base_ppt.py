#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Analizar la estructura de la PPT base para identificar elementos
"""

import os
from pptx import Presentation
from pptx.util import Inches, Cm, Pt

def emu_to_cm(emu):
    """Convierte EMU a cm"""
    # 1 cm = 360000 EMU
    return emu / 360000

base_file = os.path.join(os.path.dirname(__file__), 'REPORTE FACTURACIÓN BASE.pptx')

if not os.path.exists(base_file):
    print(f"ERROR: {base_file} no encontrado")
    exit(1)

prs = Presentation(base_file)
print(f"\n[OK] PPT base cargada: {len(prs.slides)} slides\n")

# Analizar cada slide
for slide_idx, slide in enumerate(prs.slides):
    print(f"{'='*80}")
    print(f"SLIDE {slide_idx} - {len(slide.shapes)} elementos")
    print(f"{'='*80}")
    
    for shape_idx, shape in enumerate(slide.shapes):
        print(f"\n[{shape_idx}] ", end="")
        
        # Convertir posiciones de EMU a cm
        left_cm = emu_to_cm(shape.left)
        top_cm = emu_to_cm(shape.top)
        width_cm = emu_to_cm(shape.width)
        height_cm = emu_to_cm(shape.height)
        
        # Tipo de shape
        if shape.shape_type == 13:  # PICTURE
            print(f"IMAGEN - {shape.name}")
            print(f"    Posición: ({left_cm:.2f}, {top_cm:.2f}) cm")
            print(f"    Tamaño: {width_cm:.2f} x {height_cm:.2f} cm")
        
        elif hasattr(shape, 'text_frame'):  # TEXTBOX o SHAPE CON TEXTO
            print(f"TEXTO - {shape.name}")
            text_preview = shape.text[:50].replace('\n', ' ')
            print(f"    Contenido: '{text_preview}'")
            print(f"    Posición: ({left_cm:.2f}, {top_cm:.2f}) cm")
            print(f"    Tamaño: {width_cm:.2f} x {height_cm:.2f} cm")
            
            if hasattr(shape, 'text_frame') and shape.text_frame.paragraphs:
                p = shape.text_frame.paragraphs[0]
                if p.runs:
                    r = p.runs[0]
                    font_name = r.font.name if r.font.name else "Default"
                    font_size = r.font.size.pt if r.font.size else "Default"
                    print(f"    Font: {font_name}, {font_size}pt")
        
        else:
            print(f"SHAPE - {shape.name} (tipo: {shape.shape_type})")
            print(f"    Posición: ({left_cm:.2f}, {top_cm:.2f}) cm")
            print(f"    Tamaño: {width_cm:.2f} x {height_cm:.2f} cm")

print(f"\n{'='*80}")
print("Análisis completado")
print(f"{'='*80}\n")
